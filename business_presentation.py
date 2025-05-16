import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime
import logging

# Set up logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class PresentationGenerator:
    def __init__(self):
        """Initialize presentation with template"""
        self.prs = Presentation()
        self.brand_colors = {
            'primary': RGBColor(0, 114, 178),    # Blue
            'secondary': RGBColor(213, 94, 0),    # Orange
            'accent1': RGBColor(0, 158, 115),     # Green
            'accent2': RGBColor(230, 159, 0),     # Yellow
            'text': RGBColor(55, 55, 55)          # Dark Gray
        }
    
    def add_title_slide(self):
        """Add title slide"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "E-Commerce Business Analysis"
        subtitle.text = f"Data-Driven Insights & Recommendations\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Style the text
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.brand_colors['primary']
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.brand_colors['text']
    
    def add_executive_summary(self):
        """Add executive summary slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Executive Summary"
        
        summary_points = [
            "Sales Growth: 15% projected increase",
            "Customer Segments: 6 distinct segments identified",
            "Regional Performance: South zone leads (38.90%)",
            "Product Categories: Electronics dominates (59.7%)",
            "Growth Opportunities: Strong cross-selling potential"
        ]
        
        content_text = content.text_frame
        for point in summary_points:
            p = content_text.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(18)
            p.font.color.rgb = self.brand_colors['text']
    
    def add_sales_performance(self):
        """Add sales performance slide"""
        slide_layout = self.prs.slide_layouts[2]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Sales Performance Analysis"
        
        # Add sales trend chart
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        slide.shapes.add_picture("sales_trend.png", left, top, width, height)
    
    def add_customer_segments(self):
        """Add customer segments slide"""
        slide_layout = self.prs.slide_layouts[2]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Customer Behavior Analysis"
        
        segments = {
            "Recent Customers": 29.4,
            "Loyal Customers": 23.5,
            "Big Spenders": 17.6,
            "At Risk": 5.9
        }
        
        # Create pie chart
        plt.figure(figsize=(8, 6))
        plt.pie(segments.values(), labels=segments.keys(), autopct='%1.1f%%')
        plt.title("Customer Segments Distribution")
        plt.axis('equal')
        plt.savefig('customer_segments.png')
        plt.close()
        
        # Add chart to slide
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        slide.shapes.add_picture("customer_segments.png", left, top, width, height)
    
    def add_geographic_analysis(self):
        """Add geographic analysis slide"""
        slide_layout = self.prs.slide_layouts[2]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Geographic Analysis"
        
        regions = {
            "South Zone": 38.90,
            "North Zone": 28.45,
            "West Zone": 21.65,
            "East Zone": 11.00
        }
        
        # Create bar chart
        plt.figure(figsize=(8, 6))
        plt.bar(regions.keys(), regions.values())
        plt.title("Regional Sales Distribution")
        plt.ylabel("Sales Percentage")
        plt.xticks(rotation=45)
        plt.tight_layout()
        plt.savefig('regional_sales.png')
        plt.close()
        
        # Add chart to slide
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        slide.shapes.add_picture("regional_sales.png", left, top, width, height)
    
    def add_recommendations(self):
        """Add recommendations slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Key Recommendations"
        
        recommendations = [
            "Launch Customer Retention Program",
            "Optimize Product Categories",
            "Expand in Tier 2 Cities",
            "Implement Cross-Selling Strategy",
            "Enhance Digital Experience"
        ]
        
        content_text = content.text_frame
        for rec in recommendations:
            p = content_text.add_paragraph()
            p.text = "• " + rec
            p.font.size = Pt(18)
            p.font.color.rgb = self.brand_colors['text']
    
    def add_next_steps(self):
        """Add next steps slide"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Next Steps"
        
        steps = [
            "Implement recommendations within 30 days",
            "Monitor KPIs and adjust strategies",
            "Regular reporting and stakeholder updates",
            "Quarterly review and strategy refinement"
        ]
        
        content_text = content.text_frame
        for step in steps:
            p = content_text.add_paragraph()
            p.text = "• " + step
            p.font.size = Pt(18)
            p.font.color.rgb = self.brand_colors['text']
    
    def generate_presentation(self):
        """Generate complete presentation"""
        try:
            self.add_title_slide()
            self.add_executive_summary()
            self.add_sales_performance()
            self.add_customer_segments()
            self.add_geographic_analysis()
            self.add_recommendations()
            self.add_next_steps()
            
            # Save presentation
            self.prs.save('business_analysis_presentation.pptx')
            logger.info("Presentation generated successfully!")
            
        except Exception as e:
            logger.error(f"Error generating presentation: {str(e)}")
            raise

def main():
    """Main function to generate presentation"""
    try:
        generator = PresentationGenerator()
        generator.generate_presentation()
        
    except Exception as e:
        logger.error(f"Error in main execution: {str(e)}")
        raise

if __name__ == "__main__":
    main() 